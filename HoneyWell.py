
# coding: utf-8

# In[1]:


import time
import random
from PyOCR import *
from Summarizer import *
import subprocess
from os import listdir
from os.path import isfile, join
import os.path
import sklearn.datasets
import nltk.stem
import numpy as np
import pandas as pd
from stop_words import get_stop_words
import re
import scipy as sp
import PyPDF2
import docx
import sys
import os
import os.path
import string
import docx
import docx2txt
from bs4 import BeautifulSoup

from sklearn.feature_extraction.text import TfidfVectorizer


# In[3]:


stop_words = get_stop_words('english')
english_stemmer = nltk.stem.SnowballStemmer('english')

# In[5]:


def getPdf(file):
    data = []
    file = open(file, 'rb')
    pdfReader = PyPDF2.PdfFileReader(file)
    page_number = pdfReader.numPages
    for i in range(0, page_number):
        pageObj = pdfReader.getPage(i)
        data.append(pageObj.extractText())
    return '\n'.join(data)
    


# In[6]:


def getDocx(filename):
    my_text = docx2txt.process(filename)
    
    return np.array(my_text)
    


# In[7]:


data1 = getDocx('2.docx')


# In[8]:


def getText(filename):
    doc = docx.Document(filename)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n'.join(fullText)
    
    


# In[9]:


data4 = getPdf('Data/7.pdf')


# In[10]:


import nltk
# nltk.download('punkt')
# nltk.download()
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
from nltk.probability import FreqDist
from collections import defaultdict
from string import punctuation
from heapq import nlargest


def summarize(text, n):
    sents = sent_tokenize(text)
    #assert n <= len(sents)
    if(n>=len(sents)):
        n = random.randint(int(len(sents)/2), len(sents))
    wordSent = word_tokenize(text.lower())
    stopWords = set(stopwords.words('english')+list(punctuation))
    wordSent= [word for word in wordSent if word not in stopWords]
    freq = FreqDist(wordSent)
    ranking = defaultdict(int)
    for i, sent in enumerate(sents):
        for w in word_tokenize(sent.lower()):
            if w in freq:
                ranking[i] += freq[w]
    sentsIDX = nlargest(n, ranking, key=ranking.get)
    print('Summarise Completed')
    return [sents[j] for j in sorted(sentsIDX)]

summaryArr = summarize(data4, 20)
# summaryArr


# In[13]:


from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.cluster import KMeans
import numpy as np
def Vect(data, T):
    vectorizer = TfidfVectorizer(max_df=0.5,min_df=2,stop_words='english')
    X = vectorizer.fit_transform(summaryArr)
    km = KMeans(n_clusters = T, init = 'k-means++', max_iter = 100, n_init = 2, verbose = True)
    km.fit(X)
    np.unique(km.labels_, return_counts=True)
    text={}
    for i,cluster in enumerate(km.labels_):
        oneDocument = summaryArr[i]
        if cluster not in text.keys():
            text[cluster] = oneDocument
        else:
            text[cluster] += oneDocument
    stopWords = set(stopwords.words('english')+list(punctuation))
    keywords = {}
    counts={}
    for cluster in range(T):
        word_sent = word_tokenize(text[cluster].lower())
        word_sent=[word for word in word_sent if word not in stopWords]
        freq = FreqDist(word_sent)
        keywords[cluster] = nlargest(100, freq, key=freq.get)
        counts[cluster]=freq
        uniqueKeys={}
    for cluster in range(T):   
        other_clusters=list(set(range(T))-set([cluster]))
        keys_other_clusters=set(keywords[other_clusters[0]]).union(set(keywords[other_clusters[1]]))
        unique=set(keywords[cluster])-keys_other_clusters
        uniqueKeys[cluster]=nlargest(10, unique, key=counts[cluster].get)
    #print(uniqueKeys)
    print(" ")
    print(" ")
    print('Vect Completed')
    return [X, km, vectorizer]


Vect(data4, 10)


# In[14]:


Z = Vect(data4, 10)
X = Z[0]
km = Z[1]
vectorizer = Z[2]


# In[15]:


post = """ his document is to provide an example of a weight and balance document for an aircraft type certificate application in accordance with CS-LSA. The document can be used even if the applicant does not own a DOA. It does not substitute, in any of its parts, the prescriptions of Part-21 and its amendments.
This document is intended to assist applicants in applying for an LSA RTC/TC and therefore demonstrating compliance of the design to the requirements.
The document should not be read as a template and it should not be used as a form to fill. The content shall be checked for appropriateness and changed accordingly by the applicant.
The required information can be presented entirely in this document, or in additional documents appropriately identified and referred to.
"""


# In[16]:


new_post_vec = vectorizer.transform([post])
print(new_post_vec)
#new_post_label = km.predict(new_post_vec)[0]#We will predict it's cluster 

#similar_indices = (km.labels_ == new_post_label).nonzero()[0]


# In[17]:


# Old Method, Not Doing that anymore
similar = []
for i in similar_indices:
    print(X[i])
    dist = sp.linalg.norm((new_post_vec - X[i]).toarray())
    similar.append((dist, summaryArr[i]))


# In[18]:


# New Method
similar = []
for i in range(0,20):
    dist = sp.linalg.norm((new_post_vec - X[i]).toarray())
    similar.append((dist, summaryArr[i]))


# In[19]:


similar = sorted(similar) 
print(len(similar))
show_at_1 = similar[0]
show_at_2 = similar[1]
show_at_3 = similar[2]


# In[20]:


print("=== #1 ===")
print(show_at_1)
print()


# In[21]:


print("=== #2 ===")
print(show_at_2)
print()



# In[22]:


print("=== #3 ===")
print(show_at_3)



# In[ ]:


from urllib.request import urlopen
from bs4 import BeautifulSoup
articleURL = "http://curia.europa.eu/juris/document/document.jsf?text=&docid=139407&pageIndex=0&doclang=EN&mode=lst&dir=&occ=first&part=1&cid=52454"
def getText(url):
    page = urlopen(url).read().decode('utf8', 'ignore')
    soup = BeautifulSoup(page, 'lxml')
    text = ' '.join(map(lambda p: p.text, soup.find_all('p')))
    return text.encode('ascii', errors='replace').decode().replace("?","")
text = getText(articleURL)


# In[23]:


def process(text, n):
    text1 = summarize(text, n) 
    Z = Vect(text1, T)
    X = Z[0]
    km = Z[1]
    vectorizer = Z[2]
    new_post_vec = vectorizer.transform([post])
    similar = []
    for i in range(0, 100):
        dist = sp.linalg.norm((new_post_vec - X[i]).toarray())
        similar.append((dist, summaryArr[i]))
    similar = sorted(similar) 
    show_at_1 = similar[0]
    show_at_2 = similar[1]
    show_at_3 = similar[2]
    return [show_at_1, show_at_2, show_at_2]
        
    


# In[ ]:


def final(data, post, n = 100):
    Complete = {}
    filename, type1 = os.path.splitext(data)
    if type1 == 'pdf':
        text = getPdf(data)
        Output = process(text)
    if type1 == ('ppt' or 'pptx'):
        text = #Complete for PPT
        Output = process(text)
    
    for i in range(1, 17):
        text = getPdf(str(i)+'.pdf')
        text1 = summarize(text, n) 
        T = Vect(text1)
        X = T[0]
        km = T[1]
        vectorizer = T[2]
        Complete[str(i)+'.pdf'] = X
        
        
        
        
        
        
    


# In[335]:


text = getPdf(str(1)+'.pdf')


# In[173]:


Complete = {}
for i in range(1, 16):
    print("./Data/" + str(i)+'.pdf')
    text1 = getPdf("./Data/" + str(i)+'.pdf')
    #text = getPdf('')
    text2 = summarize(text1, 10) 
    print(text2)
    Z = Vect(text2, 10)
    X = Z[0]
    km = Z[1]
    vectorizer = Z[2]
    Complete[str(i)+'.pdf'] = X


# In[586]:


# Updated Version is Below
"""
for l in range(2, 16):
    print(l)
    X = Complete[str(l)+'.pdf']

    similar = []
    for i in range(1, 10):
        print(i)
        dist = sp.linalg.norm((new_post_vec - X[i]).toarray())
        similar.append((dist, summaryArr[i]))
"""
    


# In[169]:


similar = []
for l in range(1, 16):
    name = str(l)+'.pdf'
    X = Complete[str(l)+'.pdf']
    
    for i in range(1, 10):
        dist = sp.linalg.norm((new_post_vec - X[i]).toarray())
        similar.append((dist, name, text2[i]))

    


# In[170]:


similar = sorted(similar) 
print(len(similar))
show_at_1 = similar[0]
show_at_2 = similar[1]
show_at_3 = similar[2]


# In[171]:


print(show_at_1)


# In[172]:


print(show_at_3)


# In[38]:


from pptx import Presentation


# In[46]:


# New one at the bottom
"""
#files = [x for x in os.listdir() if x.endswith(".pptx")]
files = ['2.pptx', '1.pptx', '3.pptx', '4.pptx']


for eachfile in files:
    prs = Presentation(eachfile)
    print(eachfile)
    print("----------------------")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print ((shape.text))
"""
        


# In[238]:


Complete1 = {}
files = ['2.pptx', '1.pptx', '3.pptx', '4.pptx']

def readPpt(data):
    data = []
    prs = Presentation(eachfile)
    #print(eachfile)
    #print("----------------------")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                data.append(shape.text)
                
    return '\n'.join(data)

files = ['1.pptx', '2.pptx', '3.pptx', '4.pptx']
Text = []
for eachfile in files:
    data = readPpt(eachfile)      
    Text1 = (summarize(data, 10)) 
    Text.append(Text1)
    #print(Text)
    print(len(Text))
    Z = Vect(Text1, 10)
    X1 = Z[0]
    km = Z[1]
    vectorizer = Z[2]
    #Complete[str(i)+'.pdf'] = X
    Complete[eachfile] = X1

print([len(i) for i in Text])
#print(Text[1])
print(len(Text))


# In[237]:


print(len(Text[3]))


# In[240]:


post1 = """A Missed approach is a procedure used when an instrument approach cannot be completed to a full-stop landing
The Missed Approach Segment is that part of an instrument approach procedure between the missed approach point (MAP), the missed approach way-point (MAWP), or the point of arrival at the decision height and the specified missed approach NAVAID (navigational aid), intersection, fix, or waypoint, as appropriate, at the minimum IFR altitude
"""
new_post_vec1 = vectorizer.transform([post1])
#print(new_post_vec1)
files = ['2.pptx', '1.pptx', '3.pptx', '4.pptx']
#print(Complete1)
similar1 = []
#print(Text)
for i in Complete:
    #name = str(i)+'.pptx'
    name = i
    X1 = Complete[i]#[str(i)+'.pptx']
    #print(Text)
    for j in range(0,10):
        try:
            dist = sp.linalg.norm((new_post_vec1 - X1[j]).toarray())
            #print(dist)
            print(j)
            similar1.append((dist, name, Text[i-1][j]))
        except Exception as e:
            print(e)
            break
similar1 = sorted(similar1) 
print(len(similar1))
print(similar1)
show_at_1 = similar1[0]
show_at_2 = similar1[1]
show_at_3 = similar1[2]


# In[241]:


print(show_at_1)


# In[242]:


print(show_at_2)
"""
from jsonrpclib.SimpleJSONRPCServer import SimpleJSONRPCServer

def resolveQuery(query):
    # Add definition here 
    #text = query + "IT WORKS LOLOLOLOL!" # Replace with suitable query -> Value
    post = query 
    new_post_vec = vectorizer.transform([post])
    print(new_post_vec)
    similar = []
    gg = [i for i in Complete]
    #print(Text)
    for i in range(0, len(gg)):
        #name = str(i)+'.pptx'
        name = gg[i]
        X1 = Complete[gg[i]]#[str(i)+'.pptx']
        #print(Text)
        for j in range(0,10):
            try:
                dist = sp.linalg.norm((new_post_vec1 - X1[j]).toarray())
                #print(dist)
                print(j)
                similar.append((dist, name, Text[i][j]))
            except Exception as e:
                print(e)
                break
    similar = sorted(similar) 
    print(similar)
    return similar
        

server = SimpleJSONRPCServer(('localhost', 1006))
server.register_function(resolveQuery)
print("Start server")
server.serve_forever()"""