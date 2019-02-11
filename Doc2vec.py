
# coding: utf-8

# In[1]:


import nltk
from pptx import Presentation
from sklearn.cluster import KMeans
from heapq import nlargest
from string import punctuation
from collections import defaultdict
from nltk.probability import FreqDist
from nltk.corpus import stopwords
from nltk.tokenize import sent_tokenize, word_tokenize
import time
import random
from PyOCR import *
from Summarizer import *
import subprocess
from os import listdir
from os.path import isfile, join
import os.path
import sklearn.datasets
import nltk.stem
import numpy as np
import pandas as pd
from stop_words import get_stop_words
import re
import scipy as sp
import PyPDF2
import docx
import sys
import os
import os.path
import string
import docx
import docx2txt
from bs4 import BeautifulSoup

from sklearn.feature_extraction.text import TfidfVectorizer

stop_words = get_stop_words('english')
english_stemmer = nltk.stem.SnowballStemmer('english')

vectorizer = TfidfVectorizer(max_df=0.5, min_df=2, stop_words='english')


def getPdf(file):
    data = []
    file = open(file, 'rb')
    pdfReader = PyPDF2.PdfFileReader(file)
    page_number = pdfReader.numPages
    for i in range(0, page_number):
        pageObj = pdfReader.getPage(i)
        data.append(pageObj.extractText())
    return '\n'.join(data)


def getText(filename):
    doc = docx.Document(filename)
    fullText = []
    for para in doc.paragraphs:
        fullText.append(para.text)
    return '\n'.join(fullText)


# nltk.download('punkt')
# nltk.download()


def summarize(text, n):
    sents = sent_tokenize(text)
    # assert n <= len(sents)
    if(n >= len(sents)):
        n = random.randint(int(len(sents)/2), len(sents))
    wordSent = word_tokenize(text.lower())
    stopWords = set(stopwords.words('english')+list(punctuation))
    wordSent = [word for word in wordSent if word not in stopWords]
    freq = FreqDist(wordSent)
    ranking = defaultdict(int)
    for i, sent in enumerate(sents):
        for w in word_tokenize(sent.lower()):
            if w in freq:
                ranking[i] += freq[w]
    sentsIDX = nlargest(n, ranking, key=ranking.get)
    print('Summarise Completed')
    return [sents[j] for j in sorted(sentsIDX)]

# summaryArr

# summaryArr = #summarize(data4, 20)
# summaryArr

# In[13]:


def Vect(data, T):
    try:
        vectorizer = TfidfVectorizer(
            max_df=0.7, min_df=2, stop_words='english')
        X = vectorizer.fit_transform(data)
        km = KMeans(n_clusters=T, init='k-means++',
                    max_iter=100, n_init=2, verbose=True)
        km.fit(X)
        np.unique(km.labels_, return_counts=True)
        text = {}
        for i, cluster in enumerate(km.labels_):
            oneDocument = data[i]
            if cluster not in text.keys():
                text[cluster] = oneDocument
            else:
                text[cluster] += oneDocument
        stopWords = set(stopwords.words('english')+list(punctuation))
        keywords = {}
        counts = {}
        for cluster in range(T):
            word_sent = word_tokenize(text[cluster].lower())
            word_sent = [word for word in word_sent if word not in stopWords]
            freq = FreqDist(word_sent)
            keywords[cluster] = nlargest(100, freq, key=freq.get)
            counts[cluster] = freq
            uniqueKeys = {}
        for cluster in range(T):
            other_clusters = list(set(range(T))-set([cluster]))
            keys_other_clusters = set(keywords[other_clusters[0]]).union(
                set(keywords[other_clusters[1]]))
            unique = set(keywords[cluster])-keys_other_clusters
            uniqueKeys[cluster] = nlargest(10, unique, key=counts[cluster].get)
        # print(uniqueKeys)
        print(" ")
        print(" ")
        print('Vect Completed')
        return [X, km, vectorizer]
    except Exception as e:
        print(e)
        return [None, None, None]


Complete = {}
Vectorizers = dict()
for i in range(1, 16):
    print("./Data/" + str(i)+'.pdf')
    text1 = getPdf("./Data/" + str(i)+'.pdf')
    # text = getPdf('')
    text2 = summarize(text1, 10)
    print(text2)
    Z = Vect(text2, 10)
    X = Z[0]
    km = Z[1]
    vectorizer = Z[2]
    Complete[str(i)+'.pdf'] = X
    Vectorizers[str(i)+'.pdf'] = vectorizer


files = ['2.pptx', '1.pptx', '4.pptx']


def readPpt(eachfile):
    data = []
    prs = Presentation(eachfile)
    # print(eachfile)
    # print("----------------------")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                data.append(shape.text)
    return '\n'.join(data)


Text = []
for eachfile in files:
    data = readPpt("./Data/" + eachfile)
    Text1 = (summarize(data, 10))
    Text.append(Text1)
    # print(Text)
    print(len(Text))
    Z = Vect(Text1, 10)
    X1 = Z[0]
    km = Z[1]
    vectorizer = Z[2]
    # Complete[str(i)+'.pdf'] = X
    Complete[eachfile] = X1
    Vectorizers[eachfile] = vectorizer


post1 = """A Missed approach is a procedure used when an instrument approach cannot be completed to a full-stop landing
The Missed Approach Segment is that part of an instrument approach procedure between the missed approach point (MAP), the missed approach way-point (MAWP), or the point of arrival at the decision height and the specified missed approach NAVAID (navigational aid), intersection, fix, or waypoint, as appropriate, at the minimum IFR altitude
"""
# vectorizer = TfidfVectorizer(max_df=0.5,min_df=2,stop_words='english')
# vectorizer = TfidfVectorizer(max_df=1,min_df=0.5,stop_words='english')
# new_post_vec1 = vectorizer.fit_transform([post1])
# print(new_post_vec1.size)
# print(new_post_vec1)
# print(Complete1)
"""
similar = []
# print(Text)
gg = [i for i in Complete]
for i in range(0, len(gg)):
    try:
        # name = str(i)+'.pptx'
        name = gg[i]
        X1 = Complete[gg[i]]  # [str(i)+'.pptx']
        # print(Text)
        vectorizer = Vectorizers[gg[i]]
        new_post_vec1 = vectorizer.transform([post1])
        for j in range(0, 10):
            try:
                dist = sp.linalg.norm((new_post_vec1 - X1[j]).toarray())
                print(dist)
                print(j)
                similar.append((dist, name, Text[i][j]))
            except Exception as e:
                print(e)
                break
        similar = sorted(similar) 
        print(similar)
    except Exception as e:
        print(e)
"""
from jsonrpclib.SimpleJSONRPCServer import SimpleJSONRPCServer
import re

def resolveQuery(query):
    # Add definition here 
    # text = query + "IT WORKS LOLOLOLOL!" # Replace with suitable query -> Value
    post = query 
    similar = []
    gg = [i for i in Complete]
    # print(Text)
    for i in range(0, len(gg)):
        try:
            name = gg[i]
            X1 = Complete[gg[i]]
            # print(Text)
            vectorizer = Vectorizers[gg[i]]
            new_post_vec1 = vectorizer.transform([post])
            for j in range(0,10):
                try:
                    dist = sp.linalg.norm((new_post_vec1 - X1[j]).toarray())
                    # print(dist)
                    #print(j)
                    Text[i][j].replace('\\n', '<br>')
                    Text[i][j].replace('\n', '<br>')
                    similar.append((dist, name, Text[i][j]))
                except Exception as e:
                    print(e)
                    break
        except Exception as e:
            print(e)
    similar = sorted(similar) 
    print(similar)
    gg = list()
    for i in range(0, len(similar)):
        gg.append(dict({'result-text':''.join(re.escape(similar[i][-1]).split('\\')).replace('\n', '<br>').replace('\t', '    '), 'result-image':"asd", 'result-doc-link':'google.com', 'result-doc-name':similar[i][1], 'result-modified-date':'01-2-2019', 'result-id':"123"}))
    print(gg)
    return gg

"""
def resolveQuery(query):
    gg = list()
    for i in range(0, 1):
        gg.append(dict({'result-text':"This Works", 'result-image':"asd", 'result-doc-link':'google.com', 'result-doc-name':123, 'result-modified-date':'01-2-2019', 'result-id':"123"}))
    return gg
"""
server = SimpleJSONRPCServer(('localhost', 1006))
server.register_function(resolveQuery)
print("Start server")
server.serve_forever()
